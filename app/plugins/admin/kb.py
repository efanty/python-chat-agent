import os
import uuid
import chromadb
from flask import render_template, redirect, url_for, flash, request, jsonify, current_app
from flask_login import login_required, current_user
from .main import bp, blueprint_name
from app.extensions.init_sqlalchemy import db
from app.logger import log_error, log_admin
from app.extensions.init_loginmanager import admin_required
from app.extensions.init_csrf import csrf_protected
from app.utils.settings import get_setting_int as _get_setting_int
from app.models.llm_model import LLMModel
from app.models.knowledge_base import KnowledgeBase
from app.models.knowledge_base_file import KnowledgeBaseFile
from app.utils.plugin_utils import require_plugin
from chromadb.utils.embedding_functions import OpenAIEmbeddingFunction
try:
    import docx  # 需要安装: pip install python-docx
except ImportError:
    docx = None


# ============ Knowledge Base Management ============

@bp.route("/knowledge-bases")
@admin_required
@require_plugin(blueprint_name)
def knowledge_bases():
    page = request.args.get("page", 1, type=int)
    per_page = request.args.get("per_page", _get_setting_int("admin_per_page", 20), type=int)
    pagination = KnowledgeBase.query.order_by(KnowledgeBase.created_at.desc()).paginate(
        page=page, per_page=per_page, error_out=False
    )
    kbs = pagination.items
    embeddings = LLMModel.query.filter_by(is_active=True, model_type="embedding").all()
    return render_template("admin/knowledge_bases.html", knowledge_bases=kbs, pagination=pagination, embeddings=embeddings)


@bp.route("/knowledge-bases/add", methods=["POST"])
@admin_required
@csrf_protected
@require_plugin(blueprint_name)
def knowledge_base_add():
    data = request.form
    kb = KnowledgeBase(
        name=data.get("name"),
        description=data.get("description"),
        collection_name=data.get("collection_name"),
        embedding_model_id=data.get("embedding_model_id", type=int) or None,
        is_active=data.get("is_active") == "true",
    )
    db.session.add(kb)
    db.session.commit()
    log_admin("知识库已添加 — name=%s", kb.name)
    flash("知识库添加成功。", "success")
    return redirect(url_for("admin.knowledge_bases"))


@bp.route("/knowledge-bases/<int:kb_id>/upload", methods=["POST"])
@admin_required
@csrf_protected
@require_plugin(blueprint_name)
def knowledge_base_upload(kb_id):
    """Upload documents to a knowledge base (ChromaDB collection)."""

    kb = KnowledgeBase.query.get_or_404(kb_id)
    if not kb.is_active:
        flash("知识库已禁用，无法上传。", "danger")
        return redirect(url_for("admin.knowledge_bases"))

    # Check embedding model
    emb_model = kb.embedding_model
    if not emb_model:
        flash("请先为知识库配置 Embedding 模型。", "danger")
        return redirect(url_for("admin.knowledge_bases"))

    file = request.files.get("file")
    if not file or not file.filename:
        flash("请选择要上传的文件。", "danger")
        return redirect(url_for("admin.knowledge_bases"))

    chunk_size = request.form.get("chunk_size", 500, type=int)

    # 先将上传文件保存到磁盘
    upload_dir = os.path.join(
        current_app.config.get("UPLOAD_FOLDER", "uploads"), "kb", str(kb.id)
    )
    os.makedirs(upload_dir, exist_ok=True)
    safe_name = f"{uuid.uuid4().hex}_{file.filename}"
    saved_path = os.path.join(upload_dir, safe_name)
    file.save(saved_path)
    file_size = os.path.getsize(saved_path)

    # 从已保存的磁盘文件读取内容
    saved_path_lower = saved_path.lower()
    try:
        if saved_path_lower.endswith(".pdf"):
            from pypdf import PdfReader as _PdfReader
            with open(saved_path, "rb") as f:
                reader = _PdfReader(f)
                content = "\n".join(page.extract_text() or "" for page in reader.pages)
        elif saved_path_lower.endswith(".docx"):
            from docx import Document as _Document
            doc = _Document(saved_path)
            content = "\n".join(p.text for p in doc.paragraphs)
        elif saved_path_lower.endswith(".pptx"):
            from pptx import Presentation as _Presentation
            prs = _Presentation(saved_path)
            content = "\n".join(
                shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")
            )
        else:
            with open(saved_path, "r", encoding="utf-8", errors="replace") as f:
                content = f.read()
    except ImportError as e:
        os.remove(saved_path)
        lib_name = str(e).split("'")[1] if "'" in str(e) else "?"
        flash(f"解析库未安装，请执行 pip install {lib_name}。", "danger")
        return redirect(url_for("admin.knowledge_bases"))
    except Exception as e:
        os.remove(saved_path)
        flash(f"文件解析失败: {e}", "danger")
        return redirect(url_for("admin.knowledge_bases"))

    if not content.strip():
        os.remove(saved_path)
        flash("文件内容为空。", "warning")
        return redirect(url_for("admin.knowledge_bases"))


    # Split into chunks (by paragraphs, group up to chunk_size chars)
    paragraphs = [p.strip() for p in content.split("\n") if p.strip()]
    chunks = []
    current = ""
    for para in paragraphs:
        if len(current) + len(para) < chunk_size:
            current += para + "\n"
        else:
            if current:
                chunks.append(current.strip())
            current = para + "\n"
    if current:
        chunks.append(current.strip())

    if not chunks:
        flash("无法从文件中提取有效内容。", "warning")
        return redirect(url_for("admin.knowledge_bases"))

    # Generate embeddings via OpenAI-compatible API
    api_base = emb_model.api_base or "https://api.openai.com/v1"
    api_key = emb_model.api_key or os.getenv("OPENAI_API_KEY", "")
    model_id = emb_model.model_id

    try:
        emb_fn = OpenAIEmbeddingFunction(
            api_key=api_key,
            api_base=api_base.rstrip("/"),
            model_name=model_id,
        )
        chroma_dir = current_app.config.get("CHROMA_PERSIST_DIR", "chroma_data")
        client = chromadb.PersistentClient(path=chroma_dir)
        collection = client.get_or_create_collection(
            name=kb.collection_name,
            embedding_function=emb_fn,
        )

        # 分批处理，每批最多10个chunks（DashScope限制）
        # 从配置中读取批次大小，默认为10
        batch_size = current_app.config.get("EMBEDDING_BATCH_SIZE", 10)
        total_chunks = len(chunks)
        success_chunks = 0
        
        for i in range(0, total_chunks, batch_size):
            batch_chunks = chunks[i:i+batch_size]
            batch_ids = [f"{uuid.uuid4().hex}" for _ in batch_chunks]
            batch_metadatas = [{"source": file.filename, "chunk": i + j} for j in range(len(batch_chunks))]
            
            collection.add(
                documents=batch_chunks, 
                ids=batch_ids, 
                metadatas=batch_metadatas
            )
            success_chunks += len(batch_chunks)
            current_app.logger.info(f"已处理批次 {i//batch_size + 1}/{(total_chunks + batch_size - 1)//batch_size}, chunks: {len(batch_chunks)}")

        # 更新文档计数
        kb.document_count = (kb.document_count or 0) + success_chunks
        
        # 创建文件记录
        kb_file = KnowledgeBaseFile(
            kb_id=kb.id,
            filename=file.filename,
            filepath=saved_path,
            file_size=file_size,
            chunk_count=success_chunks,
        )
        db.session.add(kb_file)
        db.session.commit()
        
        log_admin("知识库已上传 — kb_id=%d, file=%s, chunks=%d", kb.id, file.filename, success_chunks)
        flash(f"上传成功：文件「{file.filename}」已拆分为 {success_chunks} 个段落并加入知识库。", "success")
        
    except Exception as e:
        log_error("知识库上传失败 — kb_id=%d: %s", kb.id, str(e))
        flash(f"上传失败: {e}", "danger")

    return redirect(url_for("admin.knowledge_bases"))

@bp.route("/knowledge-bases/<int:kb_id>/edit", methods=["POST"])
@admin_required
@csrf_protected
@require_plugin(blueprint_name)
def knowledge_base_edit(kb_id):
    kb = KnowledgeBase.query.get_or_404(kb_id)
    data = request.form

    kb.name = data.get("name", kb.name)
    kb.description = data.get("description", kb.description)
    kb.collection_name = data.get("collection_name", kb.collection_name)
    kb.embedding_model_id = data.get("embedding_model_id", type=int) or None
    kb.is_active = data.get("is_active") == "true"

    db.session.commit()
    log_admin("知识库已编辑 — kb_id=%d, name=%s", kb.id, kb.name)
    flash("知识库更新成功。", "success")
    return redirect(url_for("admin.knowledge_bases"))


@bp.route("/knowledge-bases/<int:kb_id>/delete", methods=["POST"])
@admin_required
@csrf_protected
@require_plugin(blueprint_name)
def knowledge_base_delete(kb_id):
    log_admin("知识库已删除 — kb_id=%d", kb_id)
    kb = KnowledgeBase.query.get_or_404(kb_id)
    db.session.delete(kb)
    db.session.commit()
    flash("知识库已删除。", "success")
    return redirect(url_for("admin.knowledge_bases"))


@bp.route("/knowledge-bases/<int:kb_id>/files/<int:file_id>/delete", methods=["POST"])
@admin_required
@csrf_protected
@require_plugin(blueprint_name)
def knowledge_base_file_delete(kb_id, file_id):
    """Delete a file from knowledge base: remove ChromaDB chunks, physical file, DB record."""

    kb = KnowledgeBase.query.get_or_404(kb_id)
    kb_file = KnowledgeBaseFile.query.get_or_404(file_id)

    if kb_file.kb_id != kb.id:
        flash("\u6587\u4ef6\u4e0e\u77e5\u8bc6\u5e93\u4e0d\u5339\u914d\u3002", "danger")
        return redirect(url_for("admin.knowledge_bases"))

    try:
        # 1. Remove chunks from ChromaDB by source metadata
        emb_model = kb.embedding_model
        if emb_model:
            api_base = emb_model.api_base or "https://api.openai.com/v1"
            api_key = emb_model.api_key or os.getenv("OPENAI_API_KEY", "")
            model_id = emb_model.model_id
            emb_fn = OpenAIEmbeddingFunction(
                api_key=api_key,
                api_base=api_base.rstrip("/"),
                model_name=model_id,
            )
            chroma_dir = current_app.config.get("CHROMA_PERSIST_DIR", "chroma_data")
            client = chromadb.PersistentClient(path=chroma_dir)
            collection = client.get_or_create_collection(
                name=kb.collection_name,
                embedding_function=emb_fn,
            )
            # Query chunks with matching source filename
            result = collection.get(where={"source": kb_file.filename})
            chunk_ids = result.get("ids", [])
            if chunk_ids:
                collection.delete(ids=chunk_ids)

        # 2. Delete physical file
        if os.path.exists(kb_file.filepath):
            os.remove(kb_file.filepath)

        # 3. Delete DB record
        chunk_count = kb_file.chunk_count
        db.session.delete(kb_file)

        # 4. Update document_count
        kb.document_count = max(0, (kb.document_count or 0) - chunk_count)
        db.session.commit()

        log_admin("\u77e5\u8bc6\u5e93\u6587\u4ef6\u5df2\u5220\u9664 \u2014 kb_id=%d, file=%s, chunks=%d", kb.id, kb_file.filename, chunk_count)
        flash(f"\u6587\u4ef6\u300c{kb_file.filename}\u300d\u5df2\u5220\u9664\uff0c\u5171\u79fb\u9664 {len(chunk_ids) if chunk_ids else 0} \u4e2a\u6bb5\u843d\u3002", "success")
    except Exception as e:
        log_error("\u77e5\u8bc6\u5e93\u6587\u4ef6\u5220\u9664\u5931\u8d25 \u2014 kb_id=%d, file_id=%d: %s", kb.id, file_id, str(e))
        flash(f"\u5220\u9664\u5931\u8d25: {e}", "danger")

    return redirect(url_for("admin.knowledge_bases"))