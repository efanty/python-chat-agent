"""to_pdf Skill — 将各种格式文件转换为 PDF。

支持的输入格式:
  - Word:  .docx, .doc
  - Excel: .xlsx, .xls
  - PPT:   .pptx, .ppt
  - HTML:  .html, .htm
  - Markdown: .md
  - 纯文本: .txt
  - 图片:  .png, .jpg, .jpeg, .gif, .bmp, .tiff

中文支持:
  自动注册系统字体（微软雅黑、宋体、黑体等），确保 PDF 中正确显示中文。
"""

import os
import json
import re
from pathlib import Path
from io import BytesIO

# ── 路径 ──────────────────────────────────────────────────────────────────

_PROJECT_ROOT = Path(__file__).resolve().parent.parent.parent


def _resolve_path(file_path: str, sandbox_dir: str = None) -> Path:
    p = Path(file_path)
    if p.is_absolute():
        return p
    if sandbox_dir:
        return Path(sandbox_dir) / file_path
    return _PROJECT_ROOT / file_path


# ── 字体注册 ──────────────────────────────────────────────────────────────

_CHINESE_FONTS_REGISTERED = False


def _register_chinese_fonts():
    """注册系统中文字体到 reportlab，确保 PDF 中文显示正常。"""
    global _CHINESE_FONTS_REGISTERED
    if _CHINESE_FONTS_REGISTERED:
        return

    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont

    # 按优先级尝试注册中文字体
    font_candidates = [
        # (注册名, 字体文件路径)
        ("MicrosoftYaHei", "C:/Windows/Fonts/msyh.ttc"),
        ("SimSun", "C:/Windows/Fonts/simsun.ttc"),
        ("SimHei", "C:/Windows/Fonts/simhei.ttf"),
        ("SimKai", "C:/Windows/Fonts/simkai.ttf"),
        ("SimFang", "C:/Windows/Fonts/simfang.ttf"),
        ("HarmonySans", "C:/Windows/Fonts/HarmonyOS_Sans_SC_Regular.ttf"),
    ]

    registered = []
    for name, path in font_candidates:
        if os.path.exists(path):
            try:
                pdfmetrics.registerFont(TTFont(name, path))
                registered.append(name)
            except Exception:
                pass

    # 注册字体别名，方便统一使用
    # 优先使用微软雅黑（如果注册成功）
    if "MicrosoftYaHei" in registered:
        pdfmetrics.registerFontFamily(
            "CJK",
            normal="MicrosoftYaHei",
            bold="MicrosoftYaHei",
            italic="MicrosoftYaHei",
            boldItalic="MicrosoftYaHei",
        )
    elif "SimSun" in registered:
        pdfmetrics.registerFontFamily(
            "CJK",
            normal="SimSun",
            bold="SimSun",
            italic="SimSun",
            boldItalic="SimSun",
        )
    elif "SimHei" in registered:
        pdfmetrics.registerFontFamily(
            "CJK",
            normal="SimHei",
            bold="SimHei",
            italic="SimHei",
            boldItalic="SimHei",
        )

    _CHINESE_FONTS_REGISTERED = True


def _get_cjk_font_name() -> str:
    """获取已注册的中文字体名称。"""
    _register_chinese_fonts()
    from reportlab.pdfbase import pdfmetrics

    try:
        # 检查 CJK 字体族是否注册成功
        pdfmetrics.getFont("CJK")
        return "CJK"
    except Exception:
        pass

    # 逐个检查已知字体
    for name in ["MicrosoftYaHei", "SimSun", "SimHei", "SimKai", "SimFang", "HarmonySans"]:
        try:
            pdfmetrics.getFont(name)
            return name
        except Exception:
            pass

    return "Helvetica"  # 回退


# ── PDF 构建工具 ──────────────────────────────────────────────────────────

def _create_pdf_document(output_path: Path, title: str = ""):
    """创建 PDF 文档对象，注册中文字体。"""
    from reportlab.lib.pagesizes import A4
    from reportlab.platypus import SimpleDocTemplate
    from reportlab.lib.units import mm

    _register_chinese_fonts()
    output_path.parent.mkdir(parents=True, exist_ok=True)

    doc = SimpleDocTemplate(
        str(output_path),
        pagesize=A4,
        topMargin=20 * mm,
        bottomMargin=20 * mm,
        leftMargin=20 * mm,
        rightMargin=20 * mm,
    )
    return doc


def _make_paragraph(text: str, style_name: str = "Normal", font_size: int = 11):
    """创建段落，自动使用中文字体。"""
    from reportlab.platypus import Paragraph
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.enums import TA_LEFT, TA_CENTER

    styles = getSampleStyleSheet()
    cjk_font = _get_cjk_font_name()

    if style_name == "Title":
        base = styles["Title"]
        return Paragraph(
            text,
            ParagraphStyle(
                "CJKTile",
                parent=base,
                fontName=cjk_font,
                fontSize=22,
                leading=30,
                alignment=TA_CENTER,
            ),
        )
    elif style_name == "Heading1":
        return Paragraph(
            text,
            ParagraphStyle(
                "CJKH1",
                parent=styles["Heading1"],
                fontName=cjk_font,
                fontSize=16,
                leading=22,
                spaceBefore=16,
                spaceAfter=8,
            ),
        )
    elif style_name == "Heading2":
        return Paragraph(
            text,
            ParagraphStyle(
                "CJKH2",
                parent=styles["Heading2"],
                fontName=cjk_font,
                fontSize=14,
                leading=20,
                spaceBefore=12,
                spaceAfter=6,
            ),
        )
    elif style_name == "Heading3":
        return Paragraph(
            text,
            ParagraphStyle(
                "CJKH3",
                parent=styles["Heading3"],
                fontName=cjk_font,
                fontSize=12,
                leading=18,
                spaceBefore=10,
                spaceAfter=4,
            ),
        )
    elif style_name == "Code":
        return Paragraph(
            text,
            ParagraphStyle(
                "CJKCode",
                fontName="Courier",
                fontSize=8,
                leading=11,
                leftIndent=12,
                spaceBefore=4,
                spaceAfter=4,
                backColor=None,
            ),
        )
    else:
        return Paragraph(
            text,
            ParagraphStyle(
                "CJKNormal",
                fontName=cjk_font,
                fontSize=font_size,
                leading=font_size * 1.5,
                spaceBefore=2,
                spaceAfter=4,
            ),
        )


def _make_table(table_data: list, headers: list = None):
    """创建表格。"""
    from reportlab.platypus import Table, TableStyle
    from reportlab.lib import colors

    cjk_font = _get_cjk_font_name()

    if headers:
        data = [headers] + table_data
    else:
        data = table_data

    # 将所有单元格内容转为字符串
    str_data = []
    for row in data:
        str_data.append([str(cell) if cell is not None else "" for cell in row])

    t = Table(str_data, repeatRows=1 if headers else 0)

    style_cmds = [
        ("FONTNAME", (0, 0), (-1, -1), cjk_font),
        ("FONTSIZE", (0, 0), (-1, -1), 9),
        ("GRID", (0, 0), (-1, -1), 0.5, colors.grey),
        ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
        ("LEFTPADDING", (0, 0), (-1, -1), 6),
        ("RIGHTPADDING", (0, 0), (-1, -1), 6),
        ("TOPPADDING", (0, 0), (-1, -1), 4),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
    ]

    if headers:
        style_cmds.extend([
            ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#4472C4")),
            ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
            ("FONTNAME", (0, 0), (-1, 0), cjk_font),
            ("FONTSIZE", (0, 0), (-1, 0), 10),
            ("ALIGN", (0, 0), (-1, 0), "CENTER"),
        ])

    t.setStyle(TableStyle(style_cmds))
    return t


# ── 各格式转换函数 ────────────────────────────────────────────────────────

def _convert_docx(file_path: Path, output_path: Path, **kwargs) -> str:
    """将 Word (.docx) 文件转换为 PDF。"""
    from docx import Document

    doc = Document(str(file_path))
    pdf_doc = _create_pdf_document(output_path, title=kwargs.get("title", ""))
    story = []

    # 添加标题
    if kwargs.get("title"):
        story.append(_make_paragraph(kwargs["title"], "Title"))
        story.append(_make_paragraph("", "Normal", 6))

    # 遍历文档段落
    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            story.append(_make_paragraph(" ", "Normal", 6))
            continue

        style_name = para.style.name if para.style else "Normal"

        if "Heading 1" in style_name:
            story.append(_make_paragraph(text, "Heading1"))
        elif "Heading 2" in style_name:
            story.append(_make_paragraph(text, "Heading2"))
        elif "Heading 3" in style_name:
            story.append(_make_paragraph(text, "Heading3"))
        elif "List" in style_name:
            story.append(_make_paragraph(f"  • {text}", "Normal", 10))
        else:
            story.append(_make_paragraph(text, "Normal", 11))

    # 遍历文档表格
    for table in doc.tables:
        story.append(_make_paragraph(" ", "Normal", 6))
        table_data = []
        for row in table.rows:
            table_data.append([cell.text.strip() for cell in row.cells])
        if table_data:
            story.append(_make_table(table_data))
        story.append(_make_paragraph(" ", "Normal", 6))

    pdf_doc.build(story)
    return str(output_path)


def _convert_xlsx(file_path: Path, output_path: Path, **kwargs) -> str:
    """将 Excel (.xlsx) 文件转换为 PDF。"""
    import openpyxl

    wb = openpyxl.load_workbook(str(file_path), data_only=True)
    pdf_doc = _create_pdf_document(output_path, title=kwargs.get("title", ""))
    story = []

    if kwargs.get("title"):
        story.append(_make_paragraph(kwargs["title"], "Title"))
        story.append(_make_paragraph("", "Normal", 6))

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]

        # Sheet 标题
        story.append(_make_paragraph(f"📊 {sheet_name}", "Heading1"))

        # 读取所有数据
        rows_data = []
        for row in ws.iter_rows(values_only=True):
            row_values = [str(cell) if cell is not None else "" for cell in row]
            # 跳过全空行
            if any(v.strip() for v in row_values):
                rows_data.append(row_values)

        if rows_data:
            # 第一行作为表头
            headers = rows_data[0]
            data = rows_data[1:] if len(rows_data) > 1 else []
            if data:
                story.append(_make_table(data, headers))
            else:
                story.append(_make_table([headers]))
        else:
            story.append(_make_paragraph("（空工作表）", "Normal", 10))

        story.append(_make_paragraph(" ", "Normal", 6))

    wb.close()
    pdf_doc.build(story)
    return str(output_path)


def _convert_pptx(file_path: Path, output_path: Path, **kwargs) -> str:
    """将 PowerPoint (.pptx) 文件转换为 PDF。"""
    from pptx import Presentation

    prs = Presentation(str(file_path))
    pdf_doc = _create_pdf_document(output_path, title=kwargs.get("title", ""))
    story = []

    if kwargs.get("title"):
        story.append(_make_paragraph(kwargs["title"], "Title"))
        story.append(_make_paragraph("", "Normal", 6))

    for slide_num, slide in enumerate(prs.slides, 1):
        # 幻灯片标题
        story.append(_make_paragraph(f"📄 幻灯片 {slide_num}", "Heading2"))

        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if not text:
                        continue
                    # 检测是否像标题（字体较大或加粗）
                    if para.runs and para.runs[0].font.size:
                        font_size = para.runs[0].font.size.pt
                        if font_size >= 18:
                            story.append(_make_paragraph(text, "Heading3"))
                        else:
                            story.append(_make_paragraph(f"  • {text}", "Normal", 10))
                    else:
                        story.append(_make_paragraph(f"  • {text}", "Normal", 10))

            if shape.has_table:
                table = shape.table
                table_data = []
                for row in table.rows:
                    table_data.append([cell.text.strip() for cell in row.cells])
                if table_data:
                    story.append(_make_paragraph("", "Normal", 4))
                    story.append(_make_table(table_data))
                    story.append(_make_paragraph("", "Normal", 4))

        story.append(_make_paragraph("", "Normal", 8))

    prs = None
    pdf_doc.build(story)
    return str(output_path)


def _convert_html(file_path: Path, output_path: Path, **kwargs) -> str:
    """将 HTML 文件转换为 PDF。

    使用 BeautifulSoup 解析 HTML，提取文本内容并重建为 PDF。
    """
    from bs4 import BeautifulSoup

    with open(str(file_path), "r", encoding="utf-8", errors="replace") as f:
        html_content = f.read()

    soup = BeautifulSoup(html_content, "html.parser")

    # 提取标题
    title = kwargs.get("title", "")
    if not title:
        title_tag = soup.find("title")
        if title_tag:
            title = title_tag.get_text(strip=True)

    pdf_doc = _create_pdf_document(output_path, title=title)
    story = []

    if title:
        story.append(_make_paragraph(title, "Title"))
        story.append(_make_paragraph("", "Normal", 6))

    # 遍历 body 中的元素（使用 find_all 避免重复处理嵌套元素）
    body = soup.find("body") or soup
    for element in body.find_all():
        # 跳过嵌套在 li/pre/table 内部的子元素，避免重复处理
        if element.find_parents(["li", "pre", "table"]):
            continue
        if element.name in ("h1", "h2", "h3", "h4", "h5", "h6"):
            text = element.get_text(strip=True)
            if text:
                level = int(element.name[1])
                if level <= 2:
                    story.append(_make_paragraph(text, f"Heading{level}"))
                else:
                    story.append(_make_paragraph(text, "Heading3"))
        elif element.name == "p":
            # 跳过嵌套在 li 中的 p（li 已处理其全部文本）
            if element.find_parent("li"):
                continue
            text = element.get_text(strip=True)
            if text:
                story.append(_make_paragraph(text, "Normal", 11))
        elif element.name == "li":
            text = element.get_text(strip=True)
            if text:
                story.append(_make_paragraph(f"  • {text}", "Normal", 10))
        elif element.name == "pre":
            text = element.get_text()
            if text.strip():
                for line in text.split("\n"):
                    if line.strip():
                        story.append(_make_paragraph(line, "Code"))
        elif element.name == "table":
            rows = element.find_all("tr")
            if rows:
                table_data = []
                for row in rows:
                    cells = row.find_all(["td", "th"])
                    table_data.append([cell.get_text(strip=True) for cell in cells])
                if table_data:
                    story.append(_make_paragraph("", "Normal", 4))
                    story.append(_make_table(table_data))
                    story.append(_make_paragraph("", "Normal", 4))
        elif element.name == "hr":
            story.append(_make_paragraph("─" * 60, "Normal", 8))

    pdf_doc.build(story)
    return str(output_path)


def _convert_markdown(file_path: Path, output_path: Path, **kwargs) -> str:
    """将 Markdown (.md) 文件转换为 PDF。"""
    import markdown

    with open(str(file_path), "r", encoding="utf-8", errors="replace") as f:
        md_content = f.read()

    # 提取标题
    title = kwargs.get("title", "")
    if not title:
        # 从第一个 # 标题提取
        title_match = re.search(r"^#\s+(.+)$", md_content, re.MULTILINE)
        if title_match:
            title = title_match.group(1).strip()

    # 将 Markdown 转为 HTML
    html_body = markdown.markdown(
        md_content,
        extensions=["extra", "codehilite", "tables", "toc"],
    )

    # 包装为完整 HTML
    html_content = f"""<!DOCTYPE html>
<html><head><meta charset="utf-8"><title>{title}</title></head>
<body>{html_body}</body></html>"""

    # 写入临时 HTML 文件，再用 HTML 转换
    temp_html = output_path.with_suffix(".html")
    with open(str(temp_html), "w", encoding="utf-8") as f:
        f.write(html_content)

    try:
        result = _convert_html(temp_html, output_path, title=title)
        return result
    finally:
        # 清理临时 HTML 文件
        if temp_html.exists():
            temp_html.unlink()


def _convert_text(file_path: Path, output_path: Path, **kwargs) -> str:
    """将纯文本 (.txt) 文件转换为 PDF。"""
    with open(str(file_path), "r", encoding="utf-8", errors="replace") as f:
        text_content = f.read()

    title = kwargs.get("title", "")
    if not title:
        title = file_path.stem

    pdf_doc = _create_pdf_document(output_path, title=title)
    story = []

    story.append(_make_paragraph(title, "Title"))
    story.append(_make_paragraph("", "Normal", 6))

    for line in text_content.split("\n"):
        if line.strip():
            story.append(_make_paragraph(line, "Normal", 11))
        else:
            story.append(_make_paragraph(" ", "Normal", 6))

    pdf_doc.build(story)
    return str(output_path)


def _convert_image(file_path: Path, output_path: Path, **kwargs) -> str:
    """将图片文件转换为 PDF（每张图片一页）。"""
    from PIL import Image
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.units import mm
    from reportlab.platypus import SimpleDocTemplate, Image as RLImage, Spacer

    _register_chinese_fonts()
    output_path.parent.mkdir(parents=True, exist_ok=True)

    doc = SimpleDocTemplate(
        str(output_path),
        pagesize=A4,
        topMargin=10 * mm,
        bottomMargin=10 * mm,
        leftMargin=10 * mm,
        rightMargin=10 * mm,
    )

    story = []

    # 支持多张图片（如果 file_path 是目录或通配符）
    img_paths = [file_path]

    for img_path in img_paths:
        if not img_path.exists():
            continue

        try:
            img = Image.open(str(img_path))
            # 获取 A4 可用尺寸
            page_width = A4[0] - 20 * mm
            page_height = A4[1] - 20 * mm

            # 等比例缩放
            img_width, img_height = img.size
            ratio = min(page_width / img_width, page_height / img_height)
            display_width = img_width * ratio
            display_height = img_height * ratio

            rl_img = RLImage(str(img_path), width=display_width, height=display_height)
            story.append(rl_img)
            story.append(Spacer(1, 6 * mm))
            img.close()
        except Exception:
            story.append(_make_paragraph(f"（无法加载图片: {img_path.name}）", "Normal", 10))

    doc.build(story)
    return str(output_path)


# ── 格式检测 ──────────────────────────────────────────────────────────────

_SUPPORTED_EXTENSIONS = {
    ".docx": ("Word", _convert_docx),
    ".doc": ("Word", _convert_docx),  # .doc 也尝试用 python-docx
    ".xlsx": ("Excel", _convert_xlsx),
    ".xls": ("Excel", _convert_xlsx),  # .xls 也尝试用 openpyxl
    ".pptx": ("PowerPoint", _convert_pptx),
    ".ppt": ("PowerPoint", _convert_pptx),
    ".html": ("HTML", _convert_html),
    ".htm": ("HTML", _convert_html),
    ".md": ("Markdown", _convert_markdown),
    ".markdown": ("Markdown", _convert_markdown),
    ".txt": ("文本", _convert_text),
    ".text": ("文本", _convert_text),
    ".png": ("图片", _convert_image),
    ".jpg": ("图片", _convert_image),
    ".jpeg": ("图片", _convert_image),
    ".gif": ("图片", _convert_image),
    ".bmp": ("图片", _convert_image),
    ".tiff": ("图片", _convert_image),
    ".tif": ("图片", _convert_image),
}


# ── 入口函数 ──────────────────────────────────────────────────────────────

def run(expression: str = "", action: str = "", **kwargs) -> str:
    """将各种格式文件转换为 PDF。

    Args:
        expression: JSON 字符串或文件路径
        action: "convert"（默认）
        **kwargs:
            file_path: 源文件路径（必填）
            output: 输出 PDF 文件名（可选，默认 输入文件名.pdf）
            title: PDF 文档标题（可选）
            sandbox_dir: 沙箱目录（由系统自动传入）

    Returns:
        JSON 字符串
    """
    action = action or "convert"

    # 解析参数
    if expression and expression.strip().startswith("{"):
        try:
            args = json.loads(expression)
            file_path = args.get("file_path", kwargs.get("file_path", ""))
            output = args.get("output", kwargs.get("output", ""))
            title = args.get("title", kwargs.get("title", ""))
        except json.JSONDecodeError:
            file_path = expression
            output = kwargs.get("output", "")
            title = kwargs.get("title", "")
    else:
        file_path = kwargs.get("file_path", "") or expression
        output = kwargs.get("output", "")
        title = kwargs.get("title", "")

    if not file_path:
        return json.dumps({
            "success": False,
            "error": "请提供源文件路径。用法: run(file_path='/path/to/file.docx')",
        }, ensure_ascii=False)

    sandbox_dir = kwargs.get("sandbox_dir")
    src_path = _resolve_path(file_path, sandbox_dir)

    if not src_path.exists():
        return json.dumps({
            "success": False,
            "error": f"文件不存在: {file_path}",
        }, ensure_ascii=False)

    # 检测格式
    ext = src_path.suffix.lower()
    if ext not in _SUPPORTED_EXTENSIONS:
        return json.dumps({
            "success": False,
            "error": f"不支持的文件格式: {ext}。支持的格式: {', '.join(sorted(_SUPPORTED_EXTENSIONS.keys()))}",
        }, ensure_ascii=False)

    fmt_name, converter = _SUPPORTED_EXTENSIONS[ext]

    # 确定输出路径
    if not output:
        output = src_path.with_suffix(".pdf").name
    out_path = _resolve_path(output, sandbox_dir)

    try:
        result_path = converter(src_path, out_path, title=title)
        return json.dumps({
            "success": True,
            "output": result_path,
            "source": str(src_path),
            "source_format": fmt_name,
            "message": f"已将 {fmt_name} 文件转换为 PDF: {result_path}",
        }, ensure_ascii=False)
    except ImportError as e:
        return json.dumps({
            "success": False,
            "error": f"缺少依赖库: {e}",
        }, ensure_ascii=False)
    except Exception as e:
        return json.dumps({
            "success": False,
            "error": f"转换失败: {e}",
        }, ensure_ascii=False)


# ── Action aliases ────────────────────────────────────────────────────────

convert = run
