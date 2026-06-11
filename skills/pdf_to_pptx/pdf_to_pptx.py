"""pdf_to_pptx Skill — 使用智谱视觉大模型将 PDF 转换为 PowerPoint 演示文稿。

工作流程:
  1. PDF → 图片 (PyMuPDF/fitz)
  2. 图片 → 结构化内容 (智谱 GLM-4V 视觉模型)
  3. 结构化内容 → .pptx 文件 (python-pptx)，每页 PDF 对应一张幻灯片
"""

import os
import re
import json
import base64
import requests
from pathlib import Path
from io import BytesIO

# ── 路径 ──────────────────────────────────────────────────────────────────

_PROJECT_ROOT = Path(__file__).resolve().parent.parent.parent
_ENV_PATH = _PROJECT_ROOT / ".env"

try:
    from dotenv import load_dotenv
    load_dotenv(_ENV_PATH)
except ImportError:
    pass

# ── API 配置 ──────────────────────────────────────────────────────────────

API_URL = "https://open.bigmodel.cn/api/paas/v4/chat/completions"
DEFAULT_MODEL = os.getenv("ZHIPU_VISION_MODEL", "glm-4v-flash")
FALLBACK_MODELS = ["glm-4v-plus", "glm-5v-turbo"]

# ── 默认提示词 ────────────────────────────────────────────────────────────

DEFAULT_PPT_PROMPT = (
    "请识别此图片中的内容，以适合做幻灯片的 Markdown 格式输出。\n"
    "要求：\n"
    "1. 第一行是幻灯片标题（用 # 标记）\n"
    "2. 后续内容用要点列表（- 或 1.）呈现\n"
    "3. 如果有表格，用 Markdown 表格格式\n"
    "4. 每张幻灯片的内容要简洁、要点化\n"
    "5. 只输出 Markdown 内容，不要额外说明"
)


# ── 核心函数 ──────────────────────────────────────────────────────────────

def _resolve_path(file_path: str, sandbox_dir: str = None) -> Path:
    p = Path(file_path)
    if p.is_absolute():
        return p
    if sandbox_dir:
        return Path(sandbox_dir) / file_path
    return _PROJECT_ROOT / file_path


def _parse_page_range(pages_spec: str, total_pages: int) -> list[int]:
    if not pages_spec or pages_spec.lower() == "all":
        return list(range(1, total_pages + 1))
    pages = []
    for part in pages_spec.split(","):
        part = part.strip()
        if "-" in part:
            start, end = part.split("-", 1)
            pages.extend(range(int(start.strip()), int(end.strip()) + 1))
        else:
            pages.append(int(part))
    return sorted(set(pages))


def _pdf_to_images(pdf_path: Path, dpi: int = 200) -> list[BytesIO]:
    """将 PDF 每页转为 PNG 图片，返回 BytesIO 列表。

    使用 PyMuPDF (fitz) 实现，无需安装 poppler。
    """
    try:
        import fitz  # PyMuPDF
    except ImportError:
        raise ImportError("需要 PyMuPDF: pip install PyMuPDF")

    doc = fitz.open(str(pdf_path))
    zoom = dpi / 72
    matrix = fitz.Matrix(zoom, zoom)
    bufs = []
    for page_num in range(len(doc)):
        page = doc[page_num]
        pix = page.get_pixmap(matrix=matrix)
        img_bytes = pix.tobytes("png")
        buf = BytesIO(img_bytes)
        buf.seek(0)
        bufs.append(buf)
    doc.close()
    return bufs


def _encode_image(image_buf: BytesIO) -> str:
    b64 = base64.b64encode(image_buf.getvalue()).decode("utf-8")
    return f"data:image/png;base64,{b64}"


def _call_vision_api(data_url: str, prompt: str, model: str, api_key: str) -> tuple[bool, str, str]:
    models_to_try = [model] + [m for m in FALLBACK_MODELS if m != model]
    last_error = ""
    for attempt_model in models_to_try:
        payload = {
            "model": attempt_model,
            "messages": [{
                "role": "user",
                "content": [
                    {"type": "image_url", "image_url": {"url": data_url}},
                    {"type": "text", "text": prompt},
                ],
            }],
        }
        headers = {"Authorization": f"Bearer {api_key}"}
        try:
            resp = requests.post(API_URL, headers=headers, json=payload, timeout=120)
            if resp.status_code == 200:
                data = resp.json()
                choices = data.get("choices", [])
                if choices:
                    text = choices[0].get("message", {}).get("content", "")
                    return True, text, attempt_model
                last_error = "API 返回了空的 choices"
            else:
                last_error = f"HTTP {resp.status_code}: {resp.text[:200]}"
        except requests.exceptions.Timeout:
            last_error = "请求超时"
        except requests.exceptions.RequestException as e:
            last_error = f"请求失败: {e}"
        except Exception as e:
            last_error = f"解析错误: {e}"
    return False, last_error, model


def _markdown_to_slide_content(markdown_text: str) -> dict:
    """将 Markdown 文本解析为幻灯片内容结构。

    Returns:
        {"title": str, "bullets": [str], "table": {"headers": [str], "rows": [[str]]} | None}
    """
    lines = markdown_text.strip().split("\n")
    title = ""
    bullets = []
    table = None
    in_table = False
    table_data = []

    for line in lines:
        stripped = line.strip()
        if not stripped:
            if in_table and table_data:
                pass  # 继续收集表格
            continue

        # 标题
        if stripped.startswith("#"):
            title = re.sub(r'^#+\s*', '', stripped)
            continue

        # 表格
        if stripped.startswith("|") and stripped.endswith("|"):
            cells = [c.strip() for c in stripped.split("|")[1:-1]]
            if all(re.match(r'^[-:\s]+$', c) for c in cells):
                continue  # 跳过分隔行
            if not in_table:
                in_table = True
                table_data = []
            table_data.append(cells)
            continue
        else:
            if in_table and table_data:
                if len(table_data) >= 2:
                    table = {
                        "headers": table_data[0],
                        "rows": table_data[1:],
                    }
                table_data = []
                in_table = False

        # 列表项
        if re.match(r'^[\-\*]\s+', stripped):
            text = re.sub(r'^[\-\*]\s+', '', stripped)
            bullets.append(text)
            continue

        if re.match(r'^\d+[\.\)]\s+', stripped):
            text = re.sub(r'^\d+[\.\)]\s+', '', stripped)
            bullets.append(text)
            continue

        # 普通文本作为要点
        if stripped:
            bullets.append(stripped)

    # 处理末尾的表格
    if in_table and table_data and len(table_data) >= 2:
        table = {
            "headers": table_data[0],
            "rows": table_data[1:],
        }

    return {"title": title, "bullets": bullets, "table": table}


def _create_pptx(slides_content: list[dict], output_path: Path) -> str:
    """根据幻灯片内容列表创建 .pptx 文件。"""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        raise ImportError("需要 python-pptx: pip install python-pptx")

    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9 宽屏
    prs.slide_height = Inches(7.5)

    for slide_data in slides_content:
        slide_layout = prs.slide_layouts[6]  # blank layout
        slide = prs.slides.add_slide(slide_layout)

        # 设置背景色（浅色渐变）
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)

        title = slide_data.get("title", "")
        bullets = slide_data.get("bullets", [])
        table = slide_data.get("table")

        # ── 标题 ──────────────────────────────────────────────────────
        if title:
            txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(1.0))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(32)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)
            p.alignment = PP_ALIGN.LEFT

        # ── 要点列表 ──────────────────────────────────────────────────
        if bullets:
            top = Inches(1.6) if title else Inches(0.6)
            txBox = slide.shapes.add_textbox(Inches(0.8), top, Inches(11.7), Inches(5.5))
            tf = txBox.text_frame
            tf.word_wrap = True

            for i, bullet in enumerate(bullets):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = bullet
                p.font.size = Pt(18)
                p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
                p.space_after = Pt(8)
                p.level = 0

                # 检测是否为子要点（以 - 或缩进开头）
                if bullet.startswith("  ") or bullet.startswith("\t"):
                    p.level = 1
                    p.font.size = Pt(16)

        # ── 表格 ──────────────────────────────────────────────────────
        if table and table["headers"] and table["rows"]:
            headers = table["headers"]
            rows = table["rows"]
            num_rows = len(rows) + 1  # +1 for header
            num_cols = len(headers)

            # 计算表格位置
            table_top = Inches(1.6) if title else Inches(0.6)
            if bullets:
                table_top = Inches(4.0)  # 有要点时表格放下面

            table_width = Inches(11.7)
            table_height = Inches(min(num_rows * 0.5, 4.0))

            try:
                pptx_table = slide.shapes.add_table(
                    num_rows, num_cols,
                    Inches(0.8), table_top,
                    table_width, table_height
                ).table

                # 设置列宽
                col_width = int(table_width / num_cols)
                for c in range(num_cols):
                    pptx_table.columns[c].width = col_width

                # 表头
                for c, header in enumerate(headers):
                    cell = pptx_table.cell(0, c)
                    cell.text = header
                    for paragraph in cell.text_frame.paragraphs:
                        paragraph.font.size = Pt(14)
                        paragraph.font.bold = True
                        paragraph.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                        paragraph.alignment = PP_ALIGN.CENTER
                    # 表头背景色
                    from pptx.oxml.ns import qn
                    tcPr = cell._tc.get_or_add_tcPr()
                    solidFill = tcPr.makeelement(qn('a:solidFill'), {})
                    srgbClr = solidFill.makeelement(qn('a:srgbClr'), {'val': '4472C4'})
                    solidFill.append(srgbClr)
                    tcPr.append(solidFill)

                # 数据行
                for r, row in enumerate(rows):
                    for c, cell_text in enumerate(row):
                        if c < num_cols:
                            cell = pptx_table.cell(r + 1, c)
                            cell.text = str(cell_text)
                            for paragraph in cell.text_frame.paragraphs:
                                paragraph.font.size = Pt(12)
                                paragraph.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
                            # 交替行背景色
                            if r % 2 == 1:
                                tcPr = cell._tc.get_or_add_tcPr()
                                solidFill = tcPr.makeelement(qn('a:solidFill'), {})
                                srgbClr = solidFill.makeelement(qn('a:srgbClr'), {'val': 'E8EDF5'})
                                solidFill.append(srgbClr)
                                tcPr.append(solidFill)
            except Exception:
                pass  # 表格添加失败时忽略

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return str(output_path)


# ── 入口函数 ──────────────────────────────────────────────────────────────

def run(expression: str = "", action: str = "", **kwargs) -> str:
    """Top-level entry point called by SkillExecutor.

    将 PDF 文件转换为 PowerPoint 演示文稿。

    Args:
        expression: JSON 字符串或文件路径
        action: "convert"（默认）
        **kwargs:
            file_path: PDF 文件路径（必填）
            output: 输出 PPT 文件名（可选）
            pages: 页码范围（可选）
            prompt: 自定义识别提示词（可选）
            model: 视觉模型名称（可选）

    Returns:
        JSON 字符串
    """
    action = action or "convert"

    # 解析参数
    if expression and expression.strip().startswith("{"):
        try:
            args = json.loads(expression)
            file_path = args.get("file_path", kwargs.get("file_path", ""))
            output = args.get("output", kwargs.get("output", ""))
            pages = args.get("pages", kwargs.get("pages", "all"))
            prompt = args.get("prompt", kwargs.get("prompt", ""))
            model = args.get("model", kwargs.get("model", DEFAULT_MODEL))
        except json.JSONDecodeError:
            file_path = expression
            output = kwargs.get("output", "")
            pages = kwargs.get("pages", "all")
            prompt = kwargs.get("prompt", "")
            model = kwargs.get("model", DEFAULT_MODEL)
    else:
        file_path = kwargs.get("file_path", "") or expression
        output = kwargs.get("output", "")
        pages = kwargs.get("pages", "all")
        prompt = kwargs.get("prompt", "")
        model = kwargs.get("model", DEFAULT_MODEL)

    if not file_path:
        return json.dumps({
            "success": False,
            "error": "请提供 PDF 文件路径。用法: run(file_path='/path/to/file.pdf')",
        }, ensure_ascii=False)

    pdf_path = _resolve_path(file_path)
    if not pdf_path.exists():
        return json.dumps({
            "success": False,
            "error": f"PDF 文件不存在: {file_path}",
        }, ensure_ascii=False)

    if pdf_path.suffix.lower() != ".pdf":
        return json.dumps({
            "success": False,
            "error": f"不是 PDF 文件: {pdf_path.name}",
        }, ensure_ascii=False)

    api_key = os.getenv("ZHIPU_API_KEY", "")
    if not api_key:
        return json.dumps({
            "success": False,
            "error": "ZHIPU_API_KEY 未设置。请从 https://open.bigmodel.cn 获取 API Key。",
        }, ensure_ascii=False)

    sandbox_dir = kwargs.get("sandbox_dir")
    if not output:
        output = pdf_path.with_suffix(".pptx").name
    output_path = _resolve_path(output, sandbox_dir)

    # ── 步骤 1: PDF 转图片 ────────────────────────────────────────────
    try:
        images = _pdf_to_images(pdf_path)
    except ImportError as e:
        return json.dumps({"success": False, "error": str(e)}, ensure_ascii=False)
    except Exception as e:
        return json.dumps({
            "success": False,
            "error": f"PDF 转图片失败: {e}",

        }, ensure_ascii=False)

    total_pages = len(images)
    target_pages = _parse_page_range(pages, total_pages)

    # ── 步骤 2: 逐页识别 ──────────────────────────────────────────────
    use_prompt = prompt or DEFAULT_PPT_PROMPT
    slides_content = []
    errors = []

    for page_num in target_pages:
        if page_num < 1 or page_num > total_pages:
            errors.append(f"第 {page_num} 页不存在（共 {total_pages} 页）")
            continue

        img_buf = images[page_num - 1]
        data_url = _encode_image(img_buf)

        success, result_text, model_used = _call_vision_api(data_url, use_prompt, model, api_key)
        if not success:
            errors.append(f"第 {page_num} 页识别失败: {result_text}")
            continue

        slide_data = _markdown_to_slide_content(result_text)
        slide_data["page"] = page_num
        slide_data["model_used"] = model_used
        slides_content.append(slide_data)

    if not slides_content:
        return json.dumps({
            "success": False,
            "error": "未能从 PDF 中识别出任何内容",
            "details": errors,
            "total_pages": total_pages,
        }, ensure_ascii=False)

    # ── 步骤 3: 重建 PPT ──────────────────────────────────────────────
    try:
        pptx_path = _create_pptx(slides_content, output_path)
    except Exception as e:
        return json.dumps({
            "success": False,
            "error": f"创建 PPT 文件失败: {e}",
        }, ensure_ascii=False)

    result = {
        "success": True,
        "output": pptx_path,
        "total_pages": total_pages,
        "processed_pages": len(target_pages),
        "slides_created": len(slides_content),
        "model_used": model,
    }
    if errors:
        result["warnings"] = errors

    return json.dumps(result, ensure_ascii=False)


# ── Action aliases ────────────────────────────────────────────────────────

convert = run
